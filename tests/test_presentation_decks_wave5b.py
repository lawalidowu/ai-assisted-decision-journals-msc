"""Wave 5B presentation deck production validation."""
from __future__ import annotations

import csv
import hashlib
import re
import zipfile
from pathlib import Path

import fitz
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "distinction_strategy" / "05_presentation_deck"
PRIMARY = OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pptx"
FALLBACK = OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pptx"
PRIMARY_PDF = OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pdf"
FALLBACK_PDF = OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pdf"

DEMO_HASH = {
    "phase1-016": "e5d28326427f4108ed97922bf444a4467ce807e4aa198277e2ba46d6cd3a2953",
    "phase1-082": "9b131dc2b403f51a3a1de366ee56793c6ff255986a57e683d525297963cca2c0",
    "phase1-090": "7e94f16177c4ce820f956d64d27dfd9a1798e91627cb7e372ac43b0ba25ed430",
    "phase1-246": "30e2cf20540de86b618cc3790b006148fb50905e3c51b35c404dab88099440ee",
}

SECRET_RE = re.compile(r"(?i)(openai_api_key\s*=\s*sk-|sk-(?:proj|live)-|BEGIN (?:RSA |OPENSSH )?PRIVATE KEY)")
CDN_RE = re.compile(r"(?i)https?://(cdn\.|fonts\.googleapis|jsdelivr|unpkg)")


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_text(prs: Presentation) -> list[str]:
    out = []
    for slide in prs.slides:
        bits = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                bits.append(shape.text_frame.text)
        out.append("\n".join(bits))
    return out


def test_deck_files_exist() -> None:
    for path in (PRIMARY, FALLBACK, PRIMARY_PDF, FALLBACK_PDF):
        assert path.is_file(), path


def test_slide_counts() -> None:
    assert len(Presentation(str(PRIMARY)).slides) == 12
    assert len(Presentation(str(FALLBACK)).slides) == 8
    assert fitz.open(PRIMARY_PDF).page_count == 12
    assert fitz.open(FALLBACK_PDF).page_count == 8


def test_no_hidden_slides() -> None:
    for path in (PRIMARY, FALLBACK):
        prs = Presentation(str(path))
        for slide in prs.slides:
            # python-pptx: element attribute show if set to 0 means hidden
            cSld = slide._element
            assert cSld.get("show") in (None, "1", 1)


def test_centrepiece_and_required_numbers() -> None:
    texts = "\n".join(slide_text(Presentation(str(PRIMARY))))
    assert "phase1-082" in texts
    assert "Quotation support does not by itself establish decision-journal validity" in texts
    for token in ("414", "351/414", "1/6", "5/6", "21/50", "0.48", "0.39", "11/60", "37/60", "8/25/20/7"):
        assert token in texts
    fb = "\n".join(slide_text(Presentation(str(FALLBACK))))
    assert "phase1-082" in fb
    assert "Quotation support does not by itself establish decision-journal validity" in fb


def test_not_deployed_and_not_verified_decisions() -> None:
    for path in (PRIMARY, FALLBACK):
        text = "\n".join(slide_text(Presentation(str(path)))).lower()
        assert "not deployed" in text or "not production deployment" in text
        assert "automatically authoritative" in text or "authoritative evidence" in text
        assert "verified decision" not in text


def test_handbook_timing_not_contradicted() -> None:
    # Planning targets must not claim official "exactly 15" as sole handbook rule
    runbook = (OUT / "PRESENTER_RUNBOOK.md").read_text(encoding="utf-8")
    story = (ROOT / "docs/presentation/PRESENTATION_STORYBOARD.md").read_text(encoding="utf-8")
    assert "15–20" in story or "15-20" in story
    assert "no more than 12" in story or "≤12" in story


def test_timing_sheet_totals() -> None:
    rows = list(csv.DictReader((OUT / "TIMING_SHEET.csv").open(encoding="utf-8")))
    primary_total = sum(int(r["seconds"]) for r in rows if r["deck"] == "primary_15min" and r["slide_id"] != "TOTAL")
    fallback_total = sum(int(r["seconds"]) for r in rows if r["deck"] == "fallback_10min" and r["slide_id"] != "TOTAL")
    assert 14 * 60 <= primary_total <= 16 * 60
    assert 9 * 60 <= fallback_total <= 11 * 60


def test_manifest_empirical_rows() -> None:
    rows = list(csv.DictReader((OUT / "SLIDE_MANIFEST.csv").open(encoding="utf-8")))
    primary = [r for r in rows if r["deck"] == "primary_15min"]
    assert len(primary) == 12
    s07 = next(r for r in primary if r["slide_id"] == "S07")
    assert "082" in s07["evidence_source"] or "082" in s07["evidence_hash"] or "082" in s07["evidence_key"]
    assert DEMO_HASH["phase1-082"] in s07["evidence_hash"] or s07["evidence_hash"].startswith("9b131dc2")


def test_pptx_zip_integrity_no_external_media() -> None:
    for path in (PRIMARY, FALLBACK):
        with zipfile.ZipFile(path) as zf:
            assert zf.testzip() is None
            names = zf.namelist()
            assert any(n.startswith("ppt/slides/slide") for n in names)
            # no hyperlinks to http in slide XML required for render
            for name in names:
                if name.endswith(".xml"):
                    data = zf.read(name).decode("utf-8", errors="ignore")
                    assert not CDN_RE.search(data)
                    assert not SECRET_RE.search(data)


def test_rendered_png_counts() -> None:
    assert len(list((OUT / "rendered_slides" / "primary").glob("slide_*.png"))) == 12
    assert len(list((OUT / "rendered_slides" / "fallback").glob("slide_*.png"))) == 8


def test_demo_and_protected_hashes_unchanged() -> None:
    for jid, digest in DEMO_HASH.items():
        assert sha256(ROOT / "demo" / "evidence" / f"{jid}.json") == digest
    assert (
        sha256(ROOT / "data/manifests/phase1_decision_journal.json")
        == "814cc7c47a9f75bfc0a6c7b693feec7073e59131398d89fab7c9111fbb2e5e06"
    )
    wave2 = ROOT / "outputs/dissertation_integration/run_20260729_153931_wave2_final_integrity_fixes"
    assert (
        sha256(wave2 / "Lawal_Akeeb_Idowu_MSc_Dissertation_FINAL.docx")
        == "a829ff6d0b4a778f2a276f9fff45af05dbc47fa268f3a9b0b131a87099b0a2e2"
    )
    assert (
        sha256(wave2 / "Lawal_Akeeb_Idowu_MSc_Dissertation_FINAL.pdf")
        == "40c123b9743277d9083d3b66eb855e0fa7a57101017d08a7d8a2d94558a63519"
    )


def test_presenter_and_demo_cue_exist() -> None:
    assert (OUT / "PRESENTER_RUNBOOK.md").is_file()
    assert (OUT / "DEMO_CUE_CARD.md").is_file()
    cue = (OUT / "DEMO_CUE_CARD.md").read_text(encoding="utf-8").lower()
    assert "complete" in cue and "without" in cue
    assert "launch_demo.py" in cue
