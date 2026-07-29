#!/usr/bin/env python3
"""Build Wave 5B examiner presentation decks with python-pptx (offline, regenerable)."""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import sys
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
PRES = Path(__file__).resolve().parent
if str(PRES) not in sys.path:
    sys.path.insert(0, str(PRES))

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from presentation_content import (
    CASE_082,
    FALLBACK_IDS,
    FALLBACK_SECONDS,
    PRIMARY_SECONDS,
    SLIDES,
)
from presentation_theme import (
    AUTO,
    AUTO_BG,
    FREEZE,
    FREEZE_BG,
    HUMAN,
    HUMAN_BG,
    INK,
    LINE,
    MACHINE,
    MACHINE_BG,
    MARGIN_L,
    MUTED,
    PANEL,
    SLIDE_H,
    SLIDE_W,
    SOURCE,
    SOURCE_BG,
    WARN,
    WARN_BG,
    WHITE,
    add_paragraphs,
    add_text,
    chip,
    filled_card,
    footer,
    notes_text,
    paint_background,
    rounded_rect,
    set_run,
    title_block,
)

OUT = ROOT / "outputs" / "distinction_strategy" / "05_presentation_deck"
STAGE_COLOURS = {
    "source": (SOURCE_BG, SOURCE),
    "auto": (AUTO_BG, AUTO),
    "machine": (MACHINE_BG, MACHINE),
    "human": (HUMAN_BG, HUMAN),
}


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def freeze_banner(slide, left=Inches(8.6), top=Inches(0.28)):
    chip(
        slide,
        left,
        top,
        Inches(4.3),
        Inches(0.38),
        "Frozen research artefact — no live model call",
        FREEZE_BG,
        FREEZE,
        size=11,
    )


def build_s01(slide, total, num):
    d = SLIDES["S01"]
    paint_background(slide)
    title_block(slide, d["title"])
    freeze_banner(slide)
    filled_card(slide, MARGIN_L, Inches(1.35), Inches(12.4), Inches(1.7), PANEL, MACHINE)
    add_text(slide, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.35), d["claim"], size=24, color=INK)
    labels = [
        ("Machine candidate", MACHINE, MACHINE_BG),
        ("Automated check", AUTO, AUTO_BG),
        ("Human review", HUMAN, HUMAN_BG),
        ("Source evidence", SOURCE, SOURCE_BG),
    ]
    x = MARGIN_L
    for lab, fg, bg in labels:
        chip(slide, x, Inches(3.4), Inches(2.9), Inches(0.55), lab, bg, fg, size=14)
        x += Inches(3.1)
    add_text(
        slide,
        MARGIN_L,
        Inches(4.3),
        Inches(12.4),
        Inches(0.8),
        "EEEM004 viva presentation · UK COVID-19 Inquiry Module 2 case study\n"
        + d["caveat"],
        size=18,
        color=MUTED,
    )
    add_text(
        slide,
        MARGIN_L,
        Inches(5.4),
        Inches(12.4),
        Inches(1.0),
        "Main question: Can an LLM-assisted workflow produce useful, auditable candidates "
        "without treating generated text as authoritative evidence?",
        size=20,
        color=INK,
    )
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s02(slide, total, num):
    d = SLIDES["S02"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    acts = [("Question", HUMAN_BG, HUMAN), ("Procedure", AUTO_BG, AUTO), ("Measure", MACHINE_BG, MACHINE), ("Recall", SOURCE_BG, SOURCE)]
    x = MARGIN_L
    for lab, bg, fg in acts:
        shape = rounded_rect(slide, x, Inches(2.0), Inches(2.85), Inches(1.2), bg, fg)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, lab, size=20, bold=True, color=fg)
        x += Inches(3.1)
    add_paragraphs(slide, MARGIN_L, Inches(3.6), Inches(12.4), Inches(2.2), d["bullets"], size=24, spacing=10)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s03(slide, total, num):
    d = SLIDES["S03"]
    paint_background(slide)
    title_block(slide, d["title"])
    filled_card(slide, MARGIN_L, Inches(1.15), Inches(12.4), Inches(1.85), SOURCE_BG, SOURCE)
    add_text(slide, Inches(0.65), Inches(1.25), Inches(12.0), Inches(0.35), "Exact dissertation aim", size=14, bold=True, color=SOURCE)
    add_text(slide, Inches(0.65), Inches(1.55), Inches(12.0), Inches(1.25), d["claim"], size=18, color=INK)
    add_text(slide, MARGIN_L, Inches(3.2), Inches(12.4), Inches(0.7), "Presentation question: " + d["question"], size=18, color=MUTED)
    x = MARGIN_L
    y = Inches(4.1)
    for i, obj in enumerate(d["objectives"]):
        chip(slide, x, y, Inches(4.0), Inches(0.55), obj, PANEL, INK, size=13)
        if i % 3 == 2:
            x = MARGIN_L
            y += Inches(0.7)
        else:
            x += Inches(4.15)
    add_text(slide, MARGIN_L, Inches(5.7), Inches(12.4), Inches(0.4), "No numbered research questions in Chapter 1 — Aim + Objectives 1–6 only.", size=16, color=MUTED)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s04(slide, total, num):
    d = SLIDES["S04"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    stages = d["stages"]
    box_w = Inches(1.55)
    gap = Inches(0.12)
    start = MARGIN_L
    for i, (label, kind) in enumerate(stages):
        bg, fg = STAGE_COLOURS[kind]
        left = start + i * (box_w + gap)
        shape = rounded_rect(slide, left, Inches(2.0), box_w, Inches(2.3), bg, fg)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, label, size=13, bold=True, color=fg)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        kind_label = {
            "source": "authoritative source",
            "auto": "processing / automated check",
            "machine": "machine generation / freeze",
            "human": "human judgement",
        }[kind]
        set_run(p2, kind_label, size=11, color=MUTED)
        if i < len(stages) - 1:
            add_text(slide, left + box_w - Inches(0.05), Inches(2.9), Inches(0.25), Inches(0.35), "→", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        MARGIN_L,
        Inches(4.7),
        Inches(12.4),
        Inches(1.2),
        "Machine generation does not become authoritative because later checks exist.\n"
        + d["caveat"],
        size=20,
        color=INK,
    )
    legend_x = MARGIN_L
    for lab, kind in [("Source", "source"), ("Machine / frozen", "machine"), ("Automated", "auto"), ("Human", "human")]:
        bg, fg = STAGE_COLOURS[kind]
        chip(slide, legend_x, Inches(6.1), Inches(2.7), Inches(0.4), lab, bg, fg, size=12)
        legend_x += Inches(2.9)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s05(slide, total, num):
    d = SLIDES["S05"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    x = MARGIN_L
    y = Inches(2.0)
    for i, layer in enumerate(d["layers"]):
        chip(slide, x, y, Inches(2.95), Inches(0.85), layer, PANEL, INK, size=16)
        x += Inches(3.15)
        if (i + 1) % 4 == 0:
            x = MARGIN_L
            y += Inches(1.15)
    add_text(
        slide,
        MARGIN_L,
        Inches(5.0),
        Inches(12.4),
        Inches(1.2),
        "Different constructs answer different questions.\n"
        "Do not collapse the portfolio into one accuracy percentage.",
        size=22,
        color=INK,
    )
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s06(slide, total, num, *, compact=False):
    d = SLIDES["S06"]
    paint_background(slide)
    subtitle = d["claim"]
    if compact:
        subtitle = "Layered evaluation — no single F1. " + d["claim"]
    title_block(slide, d["title"], subtitle)
    gx = MARGIN_L
    for group in d["groups"]:
        filled_card(slide, gx, Inches(1.55), Inches(4.0), Inches(4.7), PANEL, MACHINE)
        add_text(slide, gx + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.4), group["name"], size=16, bold=True, color=MACHINE)
        yy = Inches(2.25)
        for value, label in group["items"]:
            add_text(slide, gx + Inches(0.2), yy, Inches(3.6), Inches(0.55), value, size=28, bold=True, color=INK)
            add_text(slide, gx + Inches(0.2), yy + Inches(0.5), Inches(3.6), Inches(0.35), label, size=16, color=MUTED)
            yy += Inches(1.15)
        gx += Inches(4.2)
    add_text(slide, MARGIN_L, Inches(6.45), Inches(12.4), Inches(0.35), d["caveat"] + "  ·  Supporting: 5/10/0 triangulation · 20 clusters · 50/53 · 49/50", size=14, color=MUTED)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s07(slide, total, num):
    d = SLIDES["S07"]
    c = d["case"]
    paint_background(slide)
    title_block(slide, d["title"])
    # Source
    filled_card(slide, MARGIN_L, Inches(1.15), Inches(6.05), Inches(2.5), SOURCE_BG, SOURCE)
    add_text(slide, Inches(0.65), Inches(1.25), Inches(5.6), Inches(0.35), "SOURCE EVIDENCE", size=14, bold=True, color=SOURCE)
    add_text(slide, Inches(0.65), Inches(1.7), Inches(5.6), Inches(1.6), c["source_quote"], size=20, color=INK)
    # Candidate
    filled_card(slide, Inches(6.85), Inches(1.15), Inches(6.0), Inches(2.5), MACHINE_BG, MACHINE)
    add_text(slide, Inches(7.05), Inches(1.25), Inches(5.6), Inches(0.35), "FROZEN LLM CANDIDATE", size=14, bold=True, color=MACHINE)
    add_text(slide, Inches(7.05), Inches(1.7), Inches(5.6), Inches(1.6), c["decision"], size=20, color=INK)
    # Results row
    chip(slide, MARGIN_L, Inches(3.9), Inches(4.0), Inches(0.7), "Traceability: PASS", AUTO_BG, AUTO, size=16)
    chip(slide, Inches(4.7), Inches(3.9), Inches(4.0), Inches(0.7), "Rubric B = HIGH  (supported wording)", MACHINE_BG, MACHINE, size=15)
    chip(slide, Inches(8.95), Inches(3.9), Inches(3.9), Inches(0.7), "Rubric A = NO  (not journal-valid)", WARN_BG, WARN, size=15)
    add_text(
        slide,
        MARGIN_L,
        Inches(4.85),
        Inches(12.4),
        Inches(0.45),
        f"Flag: {c['flag']} hearing administration  ·  ID {c['id']}",
        size=16,
        color=MUTED,
    )
    filled_card(slide, MARGIN_L, Inches(5.4), Inches(12.4), Inches(1.15), FREEZE_BG, FREEZE)
    add_text(
        slide,
        Inches(0.7),
        Inches(5.55),
        Inches(11.9),
        Inches(0.85),
        c["teaching"] + "\nSupported wording  ≠  valid journal entry.",
        size=22,
        bold=True,
        color=FREEZE,
    )
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s08(slide, total, num):
    d = SLIDES["S08"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    x = MARGIN_L
    for case in d["cases"]:
        filled_card(slide, x, Inches(1.5), Inches(4.0), Inches(4.5), PANEL, HUMAN)
        add_text(slide, x + Inches(0.2), Inches(1.7), Inches(3.6), Inches(0.4), case["id"], size=18, bold=True, color=HUMAN)
        add_text(slide, x + Inches(0.2), Inches(2.25), Inches(3.6), Inches(0.7), case["label"], size=20, bold=True, color=INK)
        add_text(slide, x + Inches(0.2), Inches(3.2), Inches(3.6), Inches(2.0), case["point"], size=18, color=INK)
        x += Inches(4.2)
    add_text(slide, MARGIN_L, Inches(6.2), Inches(12.4), Inches(0.45), d["caveat"] + "  Use offline demo for full passages.", size=16, color=MUTED)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s09(slide, total, num):
    d = SLIDES["S09"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    positions = [
        (MARGIN_L, Inches(1.5)),
        (Inches(6.85), Inches(1.5)),
        (MARGIN_L, Inches(4.0)),
        (Inches(6.85), Inches(4.0)),
    ]
    colours = [MACHINE_BG, AUTO_BG, HUMAN_BG, SOURCE_BG]
    borders = [MACHINE, AUTO, HUMAN, SOURCE]
    for (title, body), (left, top), bg, bd in zip(d["blocks"], positions, colours, borders):
        filled_card(slide, left, top, Inches(5.95), Inches(2.2), bg, bd)
        add_text(slide, left + Inches(0.25), top + Inches(0.25), Inches(5.4), Inches(0.4), title, size=18, bold=True, color=bd)
        add_text(slide, left + Inches(0.25), top + Inches(0.8), Inches(5.4), Inches(1.15), body, size=18, color=INK)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s10(slide, total, num):
    d = SLIDES["S10"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    y = Inches(1.35)
    for limit, impl in d["limits"]:
        filled_card(slide, MARGIN_L, y, Inches(12.4), Inches(0.78), PANEL, LINE)
        add_text(slide, Inches(0.65), y + Inches(0.12), Inches(6.3), Inches(0.55), limit, size=18, bold=True, color=INK)
        add_text(slide, Inches(7.1), y + Inches(0.12), Inches(5.4), Inches(0.55), "→ " + impl, size=16, color=MUTED)
        y += Inches(0.88)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s11(slide, total, num):
    d = SLIDES["S11"]
    paint_background(slide)
    title_block(slide, d["title"], d["claim"])
    add_paragraphs(slide, MARGIN_L, Inches(1.6), Inches(12.4), Inches(4.5), [f"• {item}" for item in d["items"]], size=24, spacing=14)
    add_text(slide, MARGIN_L, Inches(6.1), Inches(12.4), Inches(0.4), d["caveat"], size=18, color=MUTED)
    footer(slide, num, total, d["evidence_label"])
    notes_text(slide, d["notes"])


def build_s12(slide, total, num, *, fold_limits=False):
    d = SLIDES["S12"]
    paint_background(slide)
    title_block(slide, d["title"])
    if fold_limits:
        add_text(slide, MARGIN_L, Inches(1.15), Inches(12.4), Inches(0.7), "Limits carried forward: " + " · ".join(d["limits_folded"]), size=16, color=MUTED)
        top = Inches(2.0)
    else:
        top = Inches(1.6)
    filled_card(slide, MARGIN_L, top, Inches(12.4), Inches(2.6), SOURCE_BG, SOURCE)
    add_text(slide, Inches(0.75), top + Inches(0.35), Inches(11.9), Inches(2.0), d["claim"], size=24, bold=True, color=INK)
    add_text(slide, MARGIN_L, top + Inches(3.0), Inches(12.4), Inches(0.8), "Questions\nPresentation complete without live demo · optional offline walkthrough ≤2 min after the centrepiece only if time remains.", size=18, color=MUTED)
    footer(slide, num, total, d["evidence_label"])
    notes = d["notes"]
    if fold_limits:
        notes = "Speak limitations briefly, then deliver close.\n" + notes
    notes_text(slide, notes)


BUILDERS = {
    "S01": build_s01,
    "S02": build_s02,
    "S03": build_s03,
    "S04": build_s04,
    "S05": build_s05,
    "S06": build_s06,
    "S07": build_s07,
    "S08": build_s08,
    "S09": build_s09,
    "S10": build_s10,
    "S11": build_s11,
    "S12": build_s12,
}


def build_deck(slide_ids: list[str], *, fallback: bool) -> Presentation:
    prs = blank_prs()
    total = len(slide_ids)
    for i, sid in enumerate(slide_ids, start=1):
        slide = add_blank(prs)
        if sid == "S06":
            build_s06(slide, total, i, compact=fallback)
        elif sid == "S12":
            build_s12(slide, total, i, fold_limits=fallback)
        else:
            BUILDERS[sid](slide, total, i)
    return prs


def write_package_docs(primary_path: Path, fallback_path: Path) -> None:
    OUT.mkdir(parents=True, exist_ok=True)

    # TIMING_SHEET
    with (OUT / "TIMING_SHEET.csv").open("w", encoding="utf-8", newline="") as f:
        w = csv.writer(f)
        w.writerow(["deck", "slide_id", "title", "seconds", "mm_ss"])
        for sid, sec in PRIMARY_SECONDS.items():
            w.writerow(["primary_15min", sid, SLIDES[sid]["title"], sec, f"{sec // 60}:{sec % 60:02d}"])
        w.writerow(["primary_15min", "TOTAL", "", sum(PRIMARY_SECONDS.values()), ""])
        for sid in FALLBACK_IDS:
            sec = FALLBACK_SECONDS[sid]
            w.writerow(["fallback_10min", sid, SLIDES[sid]["title"], sec, f"{sec // 60}:{sec % 60:02d}"])
        w.writerow(["fallback_10min", "TOTAL", "", sum(FALLBACK_SECONDS.values()), ""])

    # SLIDE_MANIFEST
    evidence = {
        "S01": ("handbook + freeze convention", "n/a", "n/a"),
        "S02": ("dissertation/CHAPTER_1_INTRODUCTION.md", "§§1.1–1.2", "n/a"),
        "S03": ("dissertation/CHAPTER_1_INTRODUCTION.md", "§1.3 aim", "n/a"),
        "S04": ("Ch3 + demo workflow", "governed pipeline", "n/a"),
        "S05": ("docs/examiner_evidence/EXAMINER_EVIDENCE_MAP.md", "eval layers", "n/a"),
        "S06": ("phase1_decision_journal.json + eval JSON + Audit E", "multi", "814cc7c4… / Audit E"),
        "S07": ("demo/evidence/phase1-082.json", "rubric_a=no rubric_b=high", CASE_082["hash"]),
        "S08": ("demo/evidence/phase1-{016,090,246}.json", "three cases", "see Wave 4 hashes"),
        "S09": ("Ch1.6 / Ch5", "contributions", "n/a"),
        "S10": ("Ch1.4 / Ch5 / reproducibility limits", "limitations", "n/a"),
        "S11": ("Ch5 future work", "extensions", "n/a"),
        "S12": ("closing statement", "close_quote", "n/a"),
    }
    with (OUT / "SLIDE_MANIFEST.csv").open("w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(
            f,
            fieldnames=[
                "deck",
                "slide_number",
                "slide_id",
                "title",
                "principal_claim",
                "evidence_source",
                "evidence_key",
                "evidence_hash",
                "visual_asset",
                "human_adjudicated",
                "caveat",
                "speaker_time_seconds",
                "fallback_status",
                "validation_status",
            ],
        )
        w.writeheader()
        for deck, ids, secs in (
            ("primary_15min", list(PRIMARY_SECONDS), PRIMARY_SECONDS),
            ("fallback_10min", FALLBACK_IDS, FALLBACK_SECONDS),
        ):
            for i, sid in enumerate(ids, start=1):
                src, key, digest = evidence[sid]
                w.writerow(
                    {
                        "deck": deck,
                        "slide_number": i,
                        "slide_id": sid,
                        "title": SLIDES[sid]["title"],
                        "principal_claim": SLIDES[sid]["claim"][:180],
                        "evidence_source": src,
                        "evidence_key": key,
                        "evidence_hash": digest,
                        "visual_asset": "native_shapes",
                        "human_adjudicated": "yes" if sid in {"S07", "S08", "S06"} else "mixed",
                        "caveat": SLIDES[sid].get("caveat", ""),
                        "speaker_time_seconds": secs[sid],
                        "fallback_status": "included" if sid in FALLBACK_IDS else "removed_in_fallback",
                        "validation_status": "generated",
                    }
                )

    (OUT / "PRESENTER_RUNBOOK.md").write_text(
        (ROOT / "docs/presentation/PRESENTATION_SPEAKER_NOTES.md").read_text(encoding="utf-8")
        + "\n\n## Deck files\n\n"
        + f"- Primary: `{primary_path.name}`\n"
        + f"- Fallback: `{fallback_path.name}`\n"
        + "- Optional demo after S07 only (≤2 min): see DEMO_CUE_CARD.md\n",
        encoding="utf-8",
    )

    (OUT / "DEMO_CUE_CARD.md").write_text(
        """# Live demo cue card (optional, ≤2 minutes)

The presentation is **already complete** without this demo.

## When to use
Only after **S07** (phase1-082), and only if ≥2 minutes remain inside the 15–20 window.

## Sequence
1. Say: “The slide argument is complete; this is an optional frozen local walkthrough.”
2. Launch: `python demo/launch_demo.py`
3. Expected URL form: `http://127.0.0.1:<port>/index.html` (localhost only).
4. Open **082 · No × High**.
5. Point to: source quotation · frozen candidate · automated result · human A=No B=High.
6. Return immediately to the PPTX.

## Fallbacks
- `demo/print.html` if the browser cannot load cases.
- Skip entirely if timing/technical conditions are poor.
- Recovery wording: “I’ll keep to the slide evidence; the offline package remains available.”

## Do not
- Embed the live web demo in the PPTX.
- Run any live model call.
- Open more than phase1-082 during this cue.
""",
        encoding="utf-8",
    )


def export_pdf_and_pngs(pptx_path: Path, pdf_path: Path, png_dir: Path) -> dict:
    """Export via PowerPoint COM (local)."""
    import win32com.client

    png_dir.mkdir(parents=True, exist_ok=True)
    # Clean prior PNGs
    for old in png_dir.glob("*.png"):
        old.unlink()

    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    result = {"pdf": False, "png_count": 0, "errors": []}
    try:
        abs_pptx = str(pptx_path.resolve())
        presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        try:
            # 32 = ppSaveAsPDF
            presentation.SaveAs(str(pdf_path.resolve()), 32)
            result["pdf"] = pdf_path.is_file()
            # Export each slide as PNG
            for i in range(1, presentation.Slides.Count + 1):
                out_png = png_dir / f"slide_{i:02d}.png"
                presentation.Slides(i).Export(str(out_png.resolve()), "PNG", 1920, 1080)
                if out_png.is_file():
                    result["png_count"] += 1
        finally:
            presentation.Close()
    except Exception as exc:  # noqa: BLE001
        result["errors"].append(str(exc))
    finally:
        try:
            powerpoint.Quit()
        except Exception:
            pass
    return result


def write_sums() -> None:
    lines = []
    for path in sorted(OUT.rglob("*")):
        if path.is_file() and path.name != "SHA256SUMS":
            lines.append(f"{sha256(path)}  {path.relative_to(OUT).as_posix()}")
    (OUT / "SHA256SUMS").write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--skip-render", action="store_true")
    args = parser.parse_args()

    OUT.mkdir(parents=True, exist_ok=True)
    primary_ids = [f"S{i:02d}" for i in range(1, 13)]
    primary = build_deck(primary_ids, fallback=False)
    fallback = build_deck(FALLBACK_IDS, fallback=True)

    primary_path = OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pptx"
    fallback_path = OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pptx"
    primary.save(str(primary_path))
    fallback.save(str(fallback_path))
    print(f"Wrote {primary_path}")
    print(f"Wrote {fallback_path}")

    write_package_docs(primary_path, fallback_path)

    render_log = {"created_at_utc": datetime.now(timezone.utc).isoformat(), "tool": "PowerPoint COM 16"}
    if not args.skip_render:
        render_log["primary"] = export_pdf_and_pngs(
            primary_path,
            OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pdf",
            OUT / "rendered_slides" / "primary",
        )
        render_log["fallback"] = export_pdf_and_pngs(
            fallback_path,
            OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pdf",
            OUT / "rendered_slides" / "fallback",
        )
    else:
        render_log["skipped"] = True

    (OUT / "validation" / "RENDER_LOG.json").write_text(json.dumps(render_log, indent=2) + "\n", encoding="utf-8")
    write_sums()
    print("Package complete:", OUT)


if __name__ == "__main__":
    main()
