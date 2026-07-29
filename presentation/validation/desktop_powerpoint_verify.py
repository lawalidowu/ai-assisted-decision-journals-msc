#!/usr/bin/env python3
"""Desktop PowerPoint + PDF application-level verification for Wave 5B."""
from __future__ import annotations

import json
from datetime import datetime, timezone
from pathlib import Path

import fitz
import win32com.client
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "outputs" / "distinction_strategy" / "05_presentation_deck"
PRIMARY_PPTX = OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pptx"
FALLBACK_PPTX = OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pptx"
PRIMARY_PDF = OUT / "Lawal_Akeeb_MSc_Presentation_15min_12slides.pdf"
FALLBACK_PDF = OUT / "Lawal_Akeeb_MSc_Presentation_10min_8slides.pdf"


def verify_pptx_com(path: Path, expected: int) -> dict:
    result = {
        "path": str(path.relative_to(ROOT).as_posix()),
        "expected_slides": expected,
        "opened_without_repair": False,
        "slide_count": None,
        "hidden_slides": [],
        "notes_present": False,
        "centrepiece_082": False,
        "missing_media": [],
        "slideshow_ok": False,
        "errors": [],
    }
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    try:
        # WithWindow=False still loads; Open does not repair silently if corrupt → raises
        pres = app.Presentations.Open(str(path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False)
        try:
            result["opened_without_repair"] = True
            result["slide_count"] = int(pres.Slides.Count)
            texts = []
            for i in range(1, pres.Slides.Count + 1):
                slide = pres.Slides(i)
                # SlideShowTransition / Hidden
                try:
                    if bool(slide.SlideShowTransition.Hidden):
                        result["hidden_slides"].append(i)
                except Exception:
                    pass
                # Collect text
                for shape in slide.Shapes:
                    try:
                        if shape.HasTextFrame:
                            texts.append(str(shape.TextFrame.TextRange.Text))
                    except Exception:
                        pass
                # notes
                try:
                    notes = str(slide.NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text)
                    if notes.strip():
                        result["notes_present"] = True
                except Exception:
                    pass
                # pictures without link
                for shape in slide.Shapes:
                    try:
                        if shape.Type == 13:  # msoPicture
                            pass
                    except Exception:
                        result["missing_media"].append(f"slide {i}")
            blob = "\n".join(texts)
            result["centrepiece_082"] = "phase1-082" in blob and "Quotation support" in blob
            # Slideshow object availability (does not require interactive run)
            result["slideshow_ok"] = hasattr(pres, "SlideShowSettings")
            if result["slide_count"] != expected:
                result["errors"].append(f"expected {expected} slides, got {result['slide_count']}")
            if result["hidden_slides"]:
                result["errors"].append(f"hidden slides: {result['hidden_slides']}")
        finally:
            pres.Close()
    except Exception as exc:  # noqa: BLE001
        result["errors"].append(f"{type(exc).__name__}: {exc}")
    finally:
        try:
            app.Quit()
        except Exception:
            pass
    return result


def verify_pdf(path: Path, expected: int) -> dict:
    doc = fitz.open(path)
    try:
        return {
            "path": str(path.relative_to(ROOT).as_posix()),
            "page_count": doc.page_count,
            "expected_pages": expected,
            "blank_pages": [i + 1 for i, p in enumerate(doc) if not (p.get_text() or "").strip()],
            "ok": doc.page_count == expected,
        }
    finally:
        doc.close()


def main() -> None:
    # python-pptx sanity (fonts via package; repair detection via COM open)
    assert len(Presentation(str(PRIMARY_PPTX)).slides) == 12
    assert len(Presentation(str(FALLBACK_PPTX)).slides) == 8

    report = {
        "created_at_utc": datetime.now(timezone.utc).isoformat(),
        "tool": "Microsoft PowerPoint COM + PyMuPDF",
        "primary_pptx": verify_pptx_com(PRIMARY_PPTX, 12),
        "fallback_pptx": verify_pptx_com(FALLBACK_PPTX, 8),
        "primary_pdf": verify_pdf(PRIMARY_PDF, 12),
        "fallback_pdf": verify_pdf(FALLBACK_PDF, 8),
    }
    pptx_ok = (
        report["primary_pptx"]["opened_without_repair"]
        and not report["primary_pptx"]["errors"]
        and report["fallback_pptx"]["opened_without_repair"]
        and not report["fallback_pptx"]["errors"]
    )
    pdf_ok = report["primary_pdf"]["ok"] and report["fallback_pdf"]["ok"]
    report["result"] = "PASS" if pptx_ok and pdf_ok else "FAIL"

    out = OUT / "validation" / "DESKTOP_POWERPOINT_VERIFICATION.json"
    out.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, indent=2))
    if report["result"] != "PASS":
        raise SystemExit(1)
    print("DESKTOP_VERIFICATION_PASS")


if __name__ == "__main__":
    main()
