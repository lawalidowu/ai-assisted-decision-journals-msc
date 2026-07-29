"""Neutral academic theme constants and shape helpers for Wave 5B decks."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x18, 0x20, 0x28)
MUTED = RGBColor(0x4A, 0x55, 0x60)
PAPER = RGBColor(0xF4, 0xF6, 0xF8)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC5, 0xCB, 0xD3)
MACHINE = RGBColor(0x1F, 0x4E, 0x79)
MACHINE_BG = RGBColor(0xE8, 0xEE, 0xF6)
AUTO = RGBColor(0x0F, 0x5C, 0x4C)
AUTO_BG = RGBColor(0xE6, 0xF3, 0xEF)
HUMAN = RGBColor(0x1A, 0x4D, 0x6D)
HUMAN_BG = RGBColor(0xE4, 0xEE, 0xF4)
SOURCE = RGBColor(0x3D, 0x35, 0x50)
SOURCE_BG = RGBColor(0xEE, 0xEA, 0xF4)
WARN = RGBColor(0x7A, 0x1F, 0x1F)
WARN_BG = RGBColor(0xF8, 0xE8, 0xE8)
ACCENT = RGBColor(0x0B, 0x57, 0xA4)
FREEZE = RGBColor(0x5A, 0x47, 0x10)
FREEZE_BG = RGBColor(0xFF, 0xF6, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_TITLE = "Calibri"

MARGIN_L = Inches(0.45)
MARGIN_R = Inches(0.45)
MARGIN_T = Inches(0.28)
FOOTER_Y = Inches(7.15)


def set_run(paragraph, text, *, size=22, bold=False, color=INK, font=FONT):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=22,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p, text, size=size, bold=bold, color=color, font=font)
    return box


def add_paragraphs(slide, left, top, width, height, lines, *, size=20, color=INK, bold_first=False, spacing=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        set_run(p, line, size=size, bold=(bold_first and i == 0), color=color)
    return box


def rounded_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.adjustments[0] = 0.08
    return shape


def chip(slide, left, top, width, height, label, fill, text_color=WHITE, size=14):
    shape = rounded_rect(slide, left, top, width, height, fill)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, label, size=size, bold=True, color=text_color)
    return shape


def filled_card(slide, left, top, width, height, fill, border):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.25)
    return shape


def paint_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def footer(slide, slide_no, total, evidence_label, deck_label="AI-assisted decision journaling"):
    add_text(
        slide,
        MARGIN_L,
        FOOTER_Y,
        Inches(8.5),
        Inches(0.28),
        f"{deck_label}  ·  {evidence_label}",
        size=12,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(10.6),
        FOOTER_Y,
        Inches(2.3),
        Inches(0.28),
        f"{slide_no} / {total}",
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def title_block(slide, title, subtitle=None):
    add_text(slide, MARGIN_L, MARGIN_T, Inches(12.4), Inches(0.55), title, size=30, bold=True, color=INK, font=FONT_TITLE)
    if subtitle:
        add_text(slide, MARGIN_L, Inches(0.78), Inches(12.4), Inches(0.4), subtitle, size=18, color=MUTED)


def notes_text(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    set_run(p, text, size=12, color=INK)
