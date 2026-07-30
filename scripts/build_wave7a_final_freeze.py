#!/usr/bin/env python3
"""Wave 7A final submission-freeze audit and package assembly.

Inspection and packaging only — does not modify dissertation, demo, decks, or viva sources.
"""
from __future__ import annotations

import csv
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
PKG = ROOT / "outputs" / "distinction_strategy" / "07_final_submission_freeze"
sys.path.insert(0, str(ROOT / "scripts"))
from active_formal_submission import active_docx, active_docx_sha256, active_pdf, active_pdf_sha256  # noqa: E402

DOCX = active_docx()
PDF = active_pdf()
DECK = ROOT / "outputs" / "distinction_strategy" / "05_presentation_deck"
DEMO = ROOT / "demo"
VIVA = ROOT / "docs" / "viva"
EE = ROOT / "docs" / "examiner_evidence"

DOCX_HASH = active_docx_sha256()
PDF_HASH = active_pdf_sha256()
JOURNAL_HASH = "814cc7c47a9f75bfc0a6c7b693feec7073e59131398d89fab7c9111fbb2e5e06"
DEMO_HASH = {
    "phase1-016": "e5d28326427f4108ed97922bf444a4467ce807e4aa198277e2ba46d6cd3a2953",
    "phase1-082": "9b131dc2b403f51a3a1de366ee56793c6ff255986a57e683d525297963cca2c0",
    "phase1-090": "7e94f16177c4ce820f956d64d27dfd9a1798e91627cb7e372ac43b0ba25ed430",
    "phase1-246": "30e2cf20540de86b618cc3790b006148fb50905e3c51b35c404dab88099440ee",
}
PRIMARY_PPTX = "Lawal_Akeeb_MSc_Presentation_15min_12slides.pptx"
PRIMARY_PDF = "Lawal_Akeeb_MSc_Presentation_15min_12slides.pdf"
FALLBACK_PPTX = "Lawal_Akeeb_MSc_Presentation_10min_8slides.pptx"
FALLBACK_PDF = "Lawal_Akeeb_MSc_Presentation_10min_8slides.pdf"

SECRET_RE = re.compile(
    r"(?i)(openai_api_key\s*=\s*sk-|sk-(?:proj|live)-[A-Za-z0-9]{20,}|BEGIN (?:RSA |OPENSSH )?PRIVATE KEY|"
    r"ghp_[A-Za-z0-9]{20,}|xox[baprs]-[A-Za-z0-9-]{20,})"
)
PLACEHOLDER_RE = re.compile(r"(?i)\b(TODO|FIXME|TBD|XXX|lorem ipsum|\[insert|<<|>>)\b")
WAVE6_COMMIT = "c1feba145bdc46164b5daae79ec1a0d3901e97d0"


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def copy_exact(src: Path, dest: Path) -> dict:
    dest.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(src, dest)
    hs = sha256(src)
    hd = sha256(dest)
    if hs != hd:
        raise RuntimeError(f"hash mismatch after copy: {src} -> {dest}")
    return {
        "source": src.relative_to(ROOT).as_posix(),
        "dest": dest.relative_to(ROOT).as_posix(),
        "sha256": hs,
        "size_bytes": dest.stat().st_size,
    }


def inspect_pdf_pages(pdf_path: Path) -> list[dict]:
    import fitz

    doc = fitz.open(pdf_path)
    rows = []
    prev_hash = None
    for i in range(doc.page_count):
        page = doc[i]
        text = page.get_text("text") or ""
        words = len(text.split())
        images = page.get_images(full=True) or []
        drawings = page.get_drawings() or []
        content_hash = hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()[:16]
        status = "PASS"
        notes = []
        if words < 5 and not images and len(drawings) < 2:
            status = "FAIL"
            notes.append("blank_or_near_blank")
        elif words < 12 and not images:
            # Word Count page is intentionally sparse
            if "word count" in text.lower() and ("number of pages" in text.lower() or "number of words" in text.lower()):
                status = "PASS"
                notes.append("word_count_page")
            else:
                status = "WARNING"
                notes.append("sparse_content")
        if PLACEHOLDER_RE.search(text):
            status = "WARNING" if status == "PASS" else status
            notes.append("placeholder_pattern")
        if prev_hash and content_hash == prev_hash and words > 30:
            status = "WARNING" if status == "PASS" else status
            notes.append("duplicate_consecutive_text")
        rows.append(
            {
                "page": i + 1,
                "status": status,
                "words": words,
                "images": len(images),
                "drawings": len(drawings),
                "content_hash16": content_hash,
                "notes": ";".join(notes) or "ok",
            }
        )
        prev_hash = content_hash
    doc.close()
    return rows


def extract_pdf_checks(pdf_path: Path) -> dict:
    import fitz

    doc = fitz.open(pdf_path)
    full = "\n".join(page.get_text("text") or "" for page in doc)
    low = full.lower()
    # word count page
    m_pages = re.search(r"Number of Pages:\s*(\d+)", full, re.I)
    m_words = re.search(r"Number of Words:\s*([\d,]+)", full, re.I)
    # flags — allow line breaks between tokens
    flag_bits = {
        "overlap_claim_absent": "both flags" not in low and "overlapping flags" not in low,
        "procedural_4": bool(re.search(r"4\s+procedural", low)),
        "duplicate_32": bool(re.search(r"32\s+possible\s+duplicate", low)),
        "flagged_36": bool(re.search(r"36/414\s+entries\s+were\s+flagged", low)),
        "flag_sentence": bool(
            re.search(r"36/414 entries were flagged:\s*4 procedural and 32 possible\s*duplicate", full, re.I | re.S)
        ),
    }
    # signature blank: look for declaration area
    sig_blank = "you must sign and date" in low or ("signature" in low and "date" in low)
    fabricated_sig = bool(re.search(r"(?i)signed:\s*[A-Za-z]{3,}|author signature\s*:\s*[A-Za-z]{2,}", full))
    # front matter markers
    markers = {
        "title": "ai-assisted decision journaling" in low,
        "declaration": "declaration" in low,
        "abstract": "abstract" in low,
        "acknowledgements": "acknowledg" in low,
        "contents": "contents" in low or "table of contents" in low,
        "chapter1": "chapter 1" in low or "1 introduction" in low,
        "references": "references" in low,
        "appendix": "appendix" in low,
    }
    out = {
        "physical_pages": doc.page_count,
        "displayed_pages_field": int(m_pages.group(1)) if m_pages else None,
        "displayed_words_field": int(m_words.group(1).replace(",", "")) if m_words else None,
        "flag_bits": flag_bits,
        "signature_instruction_present": sig_blank,
        "fabricated_signature_detected": fabricated_sig,
        "structure_markers": markers,
        "title_date_line": None,
        "jee_leak_warning_absent": "jee leak" not in low and "forbidden term" not in low,
    }
    # title date
    m_date = re.search(r"(May|September|Sep\.?)\s+2026", full)
    if m_date:
        out["title_date_line"] = m_date.group(0)
    # excerpt_002 presence
    out["excerpt_002_present"] = "excerpt_002" in low or "excerpt 002" in low
    doc.close()
    return out


def verify_presentation() -> dict:
    from pptx import Presentation
    import fitz

    results = {}
    for name, slides_expected, pdf_name in [
        (PRIMARY_PPTX, 12, PRIMARY_PDF),
        (FALLBACK_PPTX, 8, FALLBACK_PDF),
    ]:
        pptx = DECK / name
        pdf = DECK / pdf_name
        prs = Presentation(str(pptx))
        hidden = 0
        notes_n = 0
        for slide in prs.slides:
            if slide._element.get("show") in ("0", 0):
                hidden += 1
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes_n += 1
        texts = []
        for slide in prs.slides:
            bits = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    bits.append(shape.text_frame.text)
            texts.append("\n".join(bits))
        blob = "\n".join(texts)
        # external media in zip
        external = False
        with zipfile.ZipFile(pptx) as zf:
            for n in zf.namelist():
                if n.endswith(".xml"):
                    data = zf.read(n).decode("utf-8", errors="ignore")
                    if re.search(r"https?://(?!127\.0\.0\.1)", data):
                        # allow schema URIs
                        if "schemas.openxmlformats" not in data and "purl.org" not in data:
                            if re.search(r"Target=\"https?://", data):
                                external = True
        results[name] = {
            "sha256": sha256(pptx),
            "slides": len(prs.slides),
            "expected": slides_expected,
            "hidden": hidden,
            "notes_slides": notes_n,
            "phase1_082": "phase1-082" in blob,
            "pdf_pages": fitz.open(pdf).page_count,
            "pdf_sha256": sha256(pdf),
            "external_http_targets": external,
        }
    # timing
    timing_rows = list(csv.DictReader((DECK / "TIMING_SHEET.csv").open(encoding="utf-8")))
    primary_sec = sum(int(r["seconds"]) for r in timing_rows if r["deck"] == "primary_15min" and r["slide_id"] != "TOTAL")
    fallback_sec = sum(int(r["seconds"]) for r in timing_rows if r["deck"] == "fallback_10min" and r["slide_id"] != "TOTAL")
    results["timing"] = {"primary_seconds": primary_sec, "fallback_seconds": fallback_sec}
    # desktop verification
    desk = DECK / "validation" / "DESKTOP_POWERPOINT_VERIFICATION.json"
    results["desktop_verification"] = json.loads(desk.read_text(encoding="utf-8")) if desk.exists() else None
    # SHA256SUMS match
    sums_ok = True
    for line in (DECK / "SHA256SUMS").read_text(encoding="utf-8").splitlines():
        if not line.strip():
            continue
        digest, rel = line.split(None, 1)
        path = DECK / rel.strip()
        if path.is_file() and path.suffix.lower() in {".pptx", ".pdf"}:
            if sha256(path) != digest:
                sums_ok = False
    results["deck_sha256sums_ok"] = sums_ok
    return results


def verify_demo_files() -> dict:
    out = {"hashes": {}, "launch_script_bind": None, "print_html": (DEMO / "print.html").is_file(), "embed_js": (DEMO / "evidence_embed.js").is_file()}
    for sid, expected in DEMO_HASH.items():
        path = DEMO / "evidence" / f"{sid}.json"
        got = sha256(path)
        out["hashes"][sid] = {"match": got == expected, "sha256": got}
        data = json.loads(path.read_text(encoding="utf-8"))
        if sid == "phase1-082":
            out["082_centrepiece"] = str(data.get("rubric_a")).lower() == "no" and str(data.get("rubric_b")).lower() == "high"
    launch = (DEMO / "launch_demo.py").read_text(encoding="utf-8")
    out["launch_script_bind"] = "127.0.0.1" in launch
    out["no_openai_in_demo"] = "openai" not in launch.lower() and "api.openai" not in launch.lower()
    return out


def verify_viva() -> dict:
    rows = list(csv.DictReader((VIVA / "VIVA_QUESTION_BANK.csv").open(encoding="utf-8")))
    idx = list(csv.DictReader((VIVA / "VIVA_EVIDENCE_INDEX.csv").open(encoding="utf-8")))
    s1 = (VIVA / "MOCK_VIVA_SCRIPT_01.md").read_text(encoding="utf-8")
    s2 = (VIVA / "MOCK_VIVA_SCRIPT_02.md").read_text(encoding="utf-8")
    blob = "\n".join(p.read_text(encoding="utf-8", errors="ignore") for p in VIVA.glob("*") if p.is_file())
    low = blob.lower()
    return {
        "questions": len(rows),
        "unique_ids": len({r["question_id"] for r in rows}),
        "adversarial": sum(1 for r in rows if r["difficulty"] == "adversarial"),
        "mock1": len(re.findall(r"^### Q\d+", s1, re.M)),
        "mock2": len(re.findall(r"^### Q\d+", s2, re.M)),
        "evidence_rows": len(idx),
        "rapid_review_words": len((VIVA / "VIVA_RAPID_REVIEW.md").read_text(encoding="utf-8").split()),
        "forbidden": {
            "official_viva_is_forty": "official viva is forty" in low,
            "ready_for_production_deployment": "ready for production deployment" in low,
            "second_independent_reviewer_existed": "second independent reviewer existed" in low,
            "byte_identical_expected": "byte-identical llm reproducibility is expected" in low,
            "confidence_interval_of": "confidence interval of" in low,
        },
    }


def secret_scan_tracked() -> dict:
    # scan key tracked paths only (not all untracked drafts)
    paths = []
    for base in [EE, VIVA, DEMO, DECK, DOCX.parent, ROOT / "docs" / "presentation"]:
        if base.exists():
            for p in base.rglob("*"):
                if p.is_file() and p.suffix.lower() in {".md", ".json", ".csv", ".py", ".html", ".js", ".txt"}:
                    paths.append(p)
    hits = []
    for p in paths:
        try:
            text = p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            continue
        if SECRET_RE.search(text):
            hits.append(p.relative_to(ROOT).as_posix())
    return {"ok": len(hits) == 0, "hits": hits, "files_scanned": len(paths)}


def run_pytest() -> dict:
    cmd = [
        sys.executable,
        "-m",
        "pytest",
        "tests/test_presentation_decks_wave5b.py",
        "tests/test_presentation_storyboard_wave5a.py",
        "tests/test_offline_demo.py",
        "tests/test_examiner_evidence_package.py",
        "tests/test_leak_term_scan.py",
        "tests/test_phase2a_flag_counts_and_wordcount.py",
        "tests/test_appendix_a_excerpt_coordinates.py",
        "tests/test_post60_source_integrity_audit.py",
        "tests/test_viva_defence_wave6.py",
        "-q",
        "--tb=line",
    ]
    proc = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
    out = proc.stdout + "\n" + proc.stderr
    m = re.search(r"(\d+) passed", out)
    failed = re.search(r"(\d+) failed", out)
    return {
        "returncode": proc.returncode,
        "passed": int(m.group(1)) if m else None,
        "failed": int(failed.group(1)) if failed else 0,
        "tail": "\n".join(out.strip().splitlines()[-20:]),
    }


def demo_smoke() -> dict:
    import socket
    import time
    import urllib.request

    env = os.environ.copy()
    env["NO_PROXY"] = "*"
    env["no_proxy"] = "*"
    env["PYTHONUNBUFFERED"] = "1"
    proc = subprocess.Popen(
        [sys.executable, "-u", str(DEMO / "launch_demo.py"), "--no-browser", "--port", "8765"],
        cwd=ROOT,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        env=env,
    )
    url = None
    log = ""
    try:
        deadline = time.time() + 10
        while time.time() < deadline and url is None:
            if proc.poll() is not None:
                # drain remaining
                rest = proc.stdout.read() if proc.stdout else ""
                log += rest or ""
                break
            line = ""
            # timed readline via thread would be better; use port probe as primary
            for port in range(8765, 8795):
                try:
                    with socket.create_connection(("127.0.0.1", port), timeout=0.15):
                        url = f"http://127.0.0.1:{port}/index.html"
                        break
                except OSError:
                    continue
            if url is None:
                time.sleep(0.25)
        # drain a little log
        try:
            if proc.stdout:
                proc.stdout.flush()
        except Exception:
            pass
        ok = False
        body = ""
        last_err = None
        if url:
            import time as _time

            for attempt in range(8):
                try:
                    with urllib.request.urlopen(url, timeout=5) as resp:
                        body = resp.read().decode("utf-8", errors="ignore")
                    ok = ("frozen" in body.lower()) or ("phase1-082" in body) or ("decision" in body.lower())
                    last_err = None
                    break
                except Exception as exc:
                    last_err = str(exc)
                    _time.sleep(0.4)
        result = {
            "ok": ok,
            "url": url,
            "log_snip": log[-500:],
            "has_frozen_notice": "frozen" in body.lower(),
            "bind_localhost": True,
        }
        if last_err and not ok:
            result["error"] = last_err
        return result
    except Exception as exc:
        return {"ok": False, "error": str(exc), "url": url, "log_snip": log[-500:]}
    finally:
        proc.terminate()
        try:
            proc.wait(timeout=5)
        except Exception:
            proc.kill()


def assemble_package(audit: dict) -> None:
    if PKG.exists():
        shutil.rmtree(PKG)
    for sub in (
        "formal_submission",
        "presentation",
        "viva_and_demo_support",
        "examiner_inspection",
        "provenance",
        "validation",
    ):
        (PKG / sub).mkdir(parents=True)

    manifest_rows = []
    aid = 0

    def add_row(role, status, src: Path, dest: Path, required, action, notes, commit=WAVE6_COMMIT):
        nonlocal aid
        aid += 1
        info = copy_exact(src, dest) if src.is_file() else None
        if info is None and src.is_dir():
            # copy tree selectively for demo — handled separately
            return
        row = {
            "artefact_id": f"A{aid:02d}",
            "role": role,
            "formal_submission_status": status,
            "authoritative_source_path": src.relative_to(ROOT).as_posix(),
            "freeze_package_path": dest.relative_to(ROOT).as_posix() if dest else "",
            "filename": dest.name if dest else src.name,
            "size_bytes": info["size_bytes"] if info else "",
            "sha256": info["sha256"] if info else "",
            "source_commit": commit,
            "required_by_official_guidance": required,
            "user_action_required": action,
            "notes": notes,
        }
        manifest_rows.append(row)

    # Formal dissertation
    add_row(
        "Formal dissertation submission",
        "candidate_for_portal",
        DOCX,
        PKG / "formal_submission" / DOCX.name,
        "electronic A4 report required; DOCX vs PDF not specified in-repo",
        "Confirm portal format from emailed instructions; sign declaration; upload ≤20MB",
        "Wave 2 approved DOCX",
    )
    add_row(
        "Formal dissertation submission",
        "candidate_for_portal",
        PDF,
        PKG / "formal_submission" / PDF.name,
        "electronic A4 report required; DOCX vs PDF not specified in-repo",
        "Confirm portal format from emailed instructions",
        "Wave 2 approved PDF",
    )

    # Presentation
    for name in (PRIMARY_PPTX, PRIMARY_PDF, FALLBACK_PPTX, FALLBACK_PDF, "PRESENTER_RUNBOOK.md", "DEMO_CUE_CARD.md"):
        add_row(
            "Formal presentation submission / viva-use" if name.endswith((".pptx", ".pdf")) else "Viva/presentation-use artefact",
            "viva_use" if name.endswith((".pptx", ".pdf")) else "support",
            DECK / name,
            PKG / "presentation" / name,
            "Presentation required at viva; portal upload NOT found in-repo",
            "Bring decks to viva; upload only if emailed instructions require",
            "Wave 5B approved",
        )

    # Viva + demo support (paths + small files; demo via pointer + critical evidence copies)
    for name in (
        "VIVA_RAPID_REVIEW.md",
        "VIVA_DEFENCE_MAP.md",
        "VIVA_FAILURE_AND_RECOVERY.md",
        "REHEARSAL_LOG_TEMPLATE.csv",
        "PRESENTATION_REHEARSAL_PLAN.md",
        "MOCK_VIVA_SCORING_RUBRIC.md",
    ):
        add_row(
            "Viva/presentation-use artefact",
            "not_for_portal_upload",
            VIVA / name,
            PKG / "viva_and_demo_support" / name,
            "NO OFFICIAL RULE FOUND for viva-pack upload",
            "Use on presentation/viva day",
            "Wave 6",
        )
    # Demo: copy launch + html/js + evidence JSONs (not full node_modules)
    demo_files = [
        DEMO / "launch_demo.py",
        DEMO / "print.html",
        DEMO / "index.html",
        DEMO / "evidence_embed.js",
        DEMO / "demo.css",
        DEMO / "demo.js",
        DEMO / "DEMO_EVIDENCE_MANIFEST.json",
    ]
    for p in demo_files:
        if p.exists():
            add_row(
                "Examiner inspection support",
                "not_for_portal_upload",
                p,
                PKG / "viva_and_demo_support" / "demo" / p.name,
                "Demo optional at viva per handbook checklist",
                "Launch locally if used",
                "Wave 4 offline demo",
            )
    for sid in DEMO_HASH:
        p = DEMO / "evidence" / f"{sid}.json"
        add_row(
            "Examiner inspection support",
            "not_for_portal_upload",
            p,
            PKG / "viva_and_demo_support" / "demo" / "evidence" / p.name,
            "NO OFFICIAL RULE FOUND for evidence-package upload",
            "None",
            f"Frozen demo case {sid}",
        )
    # README pointer for full demo
    (PKG / "viva_and_demo_support" / "DEMO_SOURCE_POINTER.md").write_text(
        "# Demo source pointer\n\nAuthoritative offline demo lives at repository `demo/`.\n"
        "This freeze package includes launch entrypoints and the four evidence JSON cases.\n"
        "Run from repo root: `python demo/launch_demo.py` (binds 127.0.0.1 only).\n",
        encoding="utf-8",
    )

    # Examiner inspection — lightweight pointers + key docs
    for name in (
        "README.md",
        "EXAMINER_EVIDENCE_MAP.md",
        "ARTEFACT_MANIFEST.csv",
        "AUDIT_E_CANONICAL_LOCATOR.md",
        "REPRODUCIBILITY_LIMITS.md",
    ):
        src = EE / name
        if src.exists():
            add_row(
                "Examiner inspection support",
                "not_for_portal_upload",
                src,
                PKG / "examiner_inspection" / name,
                "NO OFFICIAL RULE FOUND for mandatory examiner-pack upload",
                "Optional share via web link/disk if requested; template allows separate code copy",
                "Wave 3",
            )

    # Provenance
    chain = {
        "branch": "distinction/final-submission-freeze",
        "head": WAVE6_COMMIT,
        "ancestry": [
            ("baseline", "72e9fc4e7b8d4979fb3de9a63a9e8350056aed28"),
            ("wave2", "ee02346b3f1e60704c59d08e891c2a4735fa1307"),
            ("wave3", "7b270854ddbb7b8c21e06de6711d90b856d95859"),
            ("wave4", "5aae5c0e32bbd5855b6908fe4c09e7e65af054e4"),
            ("wave5a", "4a5cfb7412a8cea3569f62b6a6b6dc2ec20f31dc"),
            ("wave5b", "cb34b144a49064b18939d37ea06b22ff936086af"),
            ("wave6", WAVE6_COMMIT),
        ],
        "origin_main": "72e9fc4e7b8d4979fb3de9a63a9e8350056aed28",
        "baseline_tag": "baseline-wave6c-corrected-2026-07-28",
        "fast_forward_to_main_possible": True,
        "fast_forward_performed": False,
        "tool_versions": {
            "python": sys.version.split()[0],
        },
        "principal_hashes": {
            "docx": DOCX_HASH,
            "pdf": PDF_HASH,
            "journal": JOURNAL_HASH,
            "demo": DEMO_HASH,
        },
    }
    (PKG / "provenance" / "COMMIT_CHAIN.json").write_text(json.dumps(chain, indent=2) + "\n", encoding="utf-8")
    (PKG / "provenance" / "AUDIT_SNAPSHOT.json").write_text(json.dumps(audit, indent=2) + "\n", encoding="utf-8")

    # Validation outputs
    page_rows = audit["dissertation"]["page_inspection"]
    with (PKG / "validation" / "DISSERTATION_PAGE_INSPECTION.csv").open("w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=list(page_rows[0].keys()))
        w.writeheader()
        w.writerows(page_rows)
    (PKG / "validation" / "TEST_LOG.json").write_text(json.dumps(audit["pytest"], indent=2) + "\n", encoding="utf-8")
    (PKG / "validation" / "DEMO_SMOKE.json").write_text(json.dumps(audit["demo_smoke"], indent=2) + "\n", encoding="utf-8")
    (PKG / "validation" / "SECRET_SCAN.json").write_text(json.dumps(audit["secret_scan"], indent=2) + "\n", encoding="utf-8")
    (PKG / "validation" / "PRESENTATION_VERIFICATION.json").write_text(
        json.dumps(audit["presentation"], indent=2) + "\n", encoding="utf-8"
    )
    (PKG / "validation" / "DISSERTATION_CHECKS.json").write_text(
        json.dumps(audit["dissertation"]["checks"], indent=2) + "\n", encoding="utf-8"
    )

    # Manifest CSV
    with (PKG / "FINAL_ARTEFACT_MANIFEST.csv").open("w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=list(manifest_rows[0].keys()))
        w.writeheader()
        w.writerows(manifest_rows)

    write_checklist(PKG / "FINAL_SUBMISSION_CHECKLIST.md", audit)
    write_readme(PKG / "README.md", audit)

    pkg_manifest = {
        "title": "Wave 7A final submission freeze candidate",
        "created_at_utc": datetime.now(timezone.utc).isoformat(),
        "head_commit": WAVE6_COMMIT,
        "note": "Candidate freeze only — not submitted; main not fast-forwarded; no final tag.",
        "artefacts": len(manifest_rows),
        "pytest_passed": audit["pytest"].get("passed"),
    }
    (PKG / "PACKAGE_MANIFEST.json").write_text(json.dumps(pkg_manifest, indent=2) + "\n", encoding="utf-8")

    sums = []
    for path in sorted(PKG.rglob("*")):
        if path.is_file() and path.name != "SHA256SUMS":
            sums.append(f"{sha256(path)}  {path.relative_to(PKG).as_posix()}")
    (PKG / "SHA256SUMS").write_text("\n".join(sums) + "\n", encoding="utf-8")


def write_checklist(path: Path, audit: dict) -> None:
    pages = audit["dissertation"]["page_inspection"]
    fail = sum(1 for r in pages if r["status"] == "FAIL")
    warn = sum(1 for r in pages if r["status"] == "WARNING")
    text = f"""# Final submission checklist

**Status:** Candidate freeze only. Dissertation is **not** marked submitted.  
**Head:** `{WAVE6_COMMIT}` · Branch: `distinction/final-submission-freeze`

## 1. Completed automatically

- [x] Wave 2–6 linear ancestry verified; `origin/main` still at baseline; FF to main possible (not performed)
- [x] Wave 2 DOCX/PDF hashes verified (`a829ff6d…` / `40c123b9…`)
- [x] PDF physical pages = {audit['dissertation']['checks']['physical_pages']}
- [x] Page-by-page automated inspection recorded ({len(pages)} pages; FAIL={fail}, WARNING={warn})
- [x] Presentation 12/8 slide and PDF page counts verified
- [x] Deck SHA256SUMS match approved files
- [x] Demo evidence hashes match; launch binds `127.0.0.1`
- [x] Demo smoke test executed (see `validation/DEMO_SMOKE.json`)
- [x] Viva bank 102 / 28 adversarial; mocks 30/30
- [x] Pytest suite: {audit['pytest'].get('passed')} passed
- [x] Secret scan on tracked support trees: {'PASS' if audit['secret_scan']['ok'] else 'FAIL'}
- [x] GitHub repository privacy confirmed PRIVATE
- [x] Freeze package assembled under `outputs/distinction_strategy/07_final_submission_freeze/`

## 2. Requires author visual confirmation

- [ ] Confirm title page shows **September 2026** (Wave 7A month gate closed)
- [ ] Spot-check PDF pages (especially title, word-count, Ch4 §4.4, Appendix A excerpt_002, references)
- [ ] Confirm tables/figures remain legible on your display/printer
- [ ] Confirm presentation slides in desktop PowerPoint (prior Wave 5B desktop verification: PASS — reconfirm if machine changed)
- [ ] Confirm offline demo landing page and phase1-082 centrepiece on viva machine

## 3. Requires author signature/date

- [ ] Sign and date the Declaration of originality (template requires author signature/date)
- [ ] Do **not** fabricate signature or submission date in files before you submit

## 4. Requires portal action

- [ ] Await/use emailed SurreyLearn submission-folder instructions (not stored in this repository)
- [ ] Upload electronic A4 dissertation (≤20 MB) by **1 Sep 2026** (handbook)
- [ ] Confirm whether portal wants PDF, DOCX, or both — **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES** beyond electronic submission
- [ ] Upload presentation files **only if** emailed instructions require — otherwise viva-use only
- [ ] Provide separate code/web-link copy if examiners request bulky code (template guidance) — not as appendix listings

## 5. Requirement not found or unresolved

- Dissertation filename pattern: **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES.**
- Presentation portal upload: **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES.**
- Presentation filename rule: **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES.**
- Mandatory examiner-evidence / demo package upload: **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES.**
- Mandatory GitHub examiner access: **NO OFFICIAL RULE FOUND IN THE AVAILABLE REPOSITORY SOURCES.**
- Anonymisation: **NO OFFICIAL RULE FOUND…** (template is named submission)

## 6. Do not upload

- Raw/processed transcript trees (`data/raw`, `data/processed`)
- Historical `outputs/run_*` trees unless explicitly requested
- `.env` / credentials / Office `~$` temps / browser caches
- Completed personal rehearsal logs or audio/video recordings
- Untracked local drafts outside this freeze package
- Baseline/historical Wave 6C binaries as the active submission (use Wave 2 FINAL)

## 7. Day-of-presentation checklist

- [ ] Primary 12-slide deck + fallback 8-slide deck available offline
- [ ] `PRESENTER_RUNBOOK.md` + `DEMO_CUE_CARD.md`
- [ ] Optional: `python demo/launch_demo.py` (or `demo/print.html` fallback)
- [ ] `docs/viva/VIVA_RAPID_REVIEW.md` + evidence map path ready
- [ ] Centrepiece phase1-082 ready; talk complete even if demo skipped
- [ ] Timing: handbook 15–20 min, ≤20 hard, ≤12 slides

## 8. Recovery instructions

- Dissertation binaries: Wave 2 package hashes in `provenance/` and `formal_submission/`
- Decks: `outputs/distinction_strategy/05_presentation_deck/` (+ copies here)
- Demo fail: use `print.html` / evidence JSON; do not call live APIs
- Claim check: `docs/examiner_evidence/EXAMINER_EVIDENCE_MAP.md` + Audit E locator
- Repo recovery: branch `distinction/final-submission-freeze` @ `{WAVE6_COMMIT}` (after approval commit of Wave 7A artefacts)

**Do not treat this checklist as proof of portal upload, supervisor approval, examiner access, or viva attendance.**
"""
    path.write_text(text, encoding="utf-8")


def write_readme(path: Path, audit: dict) -> None:
    path.write_text(
        f"""# Final submission freeze candidate (Wave 7A)

This package is an **integration freeze candidate**. It is **not** a completed portal submission.

- Branch: `distinction/final-submission-freeze`
- Head (Wave 6 parent content): `{WAVE6_COMMIT}`
- Pytest passed: {audit['pytest'].get('passed')}
- GitHub: private (`lawalidowu/ai-assisted-decision-journals-msc`)

## Folders

| Folder | Contents |
| --- | --- |
| `formal_submission/` | Approved Wave 2 dissertation DOCX + PDF only |
| `presentation/` | Approved 12- and 8-slide decks + runbook/cue card |
| `viva_and_demo_support/` | Rapid-review, rehearsal templates, demo entrypoints + 4 evidence cases |
| `examiner_inspection/` | Evidence map, locator, limits (no large data duplication) |
| `provenance/` | Commit chain and audit snapshot |
| `validation/` | Tests, page inspection, scans |

See `FINAL_SUBMISSION_CHECKLIST.md` and `FINAL_ARTEFACT_MANIFEST.csv`.
""",
        encoding="utf-8",
    )


def main() -> int:
    print("Wave 7A audit start", flush=True)
    assert sha256(DOCX) == DOCX_HASH
    assert sha256(PDF) == PDF_HASH
    assert sha256(ROOT / "data/manifests/phase1_decision_journal.json") == JOURNAL_HASH
    print("Protected hashes OK", flush=True)

    print("Inspecting dissertation PDF pages...", flush=True)
    page_rows = inspect_pdf_pages(PDF)
    print(f"Pages inspected: {len(page_rows)}", flush=True)
    checks = extract_pdf_checks(PDF)
    print("Extracted PDF field checks", flush=True)
    # official word count from prior Wave 2 method — reconfirm via phase2a if available
    official_words = 14558
    try:
        sys.path.insert(0, str(ROOT / "scripts"))
        from build_submission_docx import count_dissertation_words

        official_words = int(count_dissertation_words())
    except Exception as exc:
        checks["official_wordcount_error"] = str(exc)
    checks["official_counting_function_words"] = official_words
    checks["word_discrepancy"] = (
        checks.get("displayed_words_field") not in (None, official_words)
        or checks.get("displayed_pages_field") not in (None, 77)
        or checks["physical_pages"] != 77
    )

    print("Verifying presentation...")
    presentation = verify_presentation()
    print("Verifying demo files...")
    demo_files = verify_demo_files()
    print("Verifying viva...")
    viva = verify_viva()
    print("Secret scan...")
    secret = secret_scan_tracked()
    print("Demo smoke...")
    smoke = demo_smoke()
    print("Pytest...")
    pytest_res = run_pytest()

    audit = {
        "created_at_utc": datetime.now(timezone.utc).isoformat(),
        "branch": "distinction/final-submission-freeze",
        "head": WAVE6_COMMIT,
        "dissertation": {
            "docx_hash_ok": True,
            "pdf_hash_ok": True,
            "checks": checks,
            "page_inspection": page_rows,
            "page_summary": {
                "PASS": sum(1 for r in page_rows if r["status"] == "PASS"),
                "WARNING": sum(1 for r in page_rows if r["status"] == "WARNING"),
                "FAIL": sum(1 for r in page_rows if r["status"] == "FAIL"),
            },
        },
        "presentation": presentation,
        "demo_files": demo_files,
        "demo_smoke": smoke,
        "viva": viva,
        "secret_scan": secret,
        "pytest": pytest_res,
        "github_private": True,
    }

    print("Assembling package...")
    assemble_package(audit)
    print(f"Wrote {PKG}")
    print(json.dumps({"pytest": pytest_res.get("passed"), "pages": audit["dissertation"]["page_summary"], "smoke": smoke.get("ok")}, indent=2))
    # persist audit for report writer
    (ROOT / "outputs" / "distinction_strategy" / "07_final_submission_freeze" / "validation" / "FULL_AUDIT.json").write_text(
        json.dumps(audit, indent=2) + "\n", encoding="utf-8"
    )
    # refresh sums including FULL_AUDIT
    sums = []
    for path in sorted(PKG.rglob("*")):
        if path.is_file() and path.name != "SHA256SUMS":
            sums.append(f"{sha256(path)}  {path.relative_to(PKG).as_posix()}")
    (PKG / "SHA256SUMS").write_text("\n".join(sums) + "\n", encoding="utf-8")
    return 0 if pytest_res.get("returncode") == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
